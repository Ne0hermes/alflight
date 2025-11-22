"""
Script pour ajouter les slides expliquant le workflow ALFlight
- Première utilisation : 10-15 minutes (configuration)
- Utilisations suivantes : 3 minutes (préparation vol)
"""

import subprocess
import sys

# Installer python-pptx si nécessaire
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
except ImportError:
    print("Installation de python-pptx...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor

def add_title_content_slide(prs, title_text, content_lines, subtitle=None, slide_index=None):
    """Ajoute une slide avec titre et contenu bullet points"""
    slide_layout = prs.slide_layouts[1]  # Layout avec titre et contenu

    if slide_index is not None:
        # Insérer à une position spécifique
        slide = prs.slides.add_slide(slide_layout)
        # Move to correct position
        xml_slides = prs.slides._sldIdLst
        slides = list(xml_slides)
        xml_slides.remove(slides[-1])
        xml_slides.insert(slide_index, slides[-1])
    else:
        slide = prs.slides.add_slide(slide_layout)

    # Titre
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(147, 22, 60)  # Bordeaux ALFlight

    # Contenu
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    # Sous-titre optionnel
    if subtitle:
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.italic = True
        p.font.color.rgb = RGBColor(107, 15, 43)
        p.space_after = Pt(20)

    # Lignes de contenu
    for i, line in enumerate(content_lines):
        if i == 0 and not subtitle:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        # Déterminer le niveau d'indentation
        if line.startswith("    "):
            p.text = line.strip()
            p.level = 2
        elif line.startswith("  "):
            p.text = line.strip()
            p.level = 1
        else:
            p.text = line
            p.level = 0

        # Tailles de police selon niveau
        if p.level == 0:
            p.font.size = Pt(18)
        elif p.level == 1:
            p.font.size = Pt(15)
        else:
            p.font.size = Pt(13)

        p.space_before = Pt(10) if p.level == 0 else Pt(6) if p.level == 1 else Pt(3)

        # Couleurs spéciales
        if "✅" in line or "SUCCESS" in line:
            p.font.color.rgb = RGBColor(16, 185, 129)  # Vert
            p.font.bold = True
        elif "⚡" in line or "RAPIDE" in line:
            p.font.color.rgb = RGBColor(59, 130, 246)  # Bleu
            p.font.bold = True

    return slide

def add_transition_slide(prs, title_text, subtitle_text=None, bg_color=None):
    """Ajoute une slide de transition"""
    slide_layout = prs.slide_layouts[5]  # Layout vide
    slide = prs.slides.add_slide(slide_layout)

    # Fond (bordeaux par défaut)
    if bg_color is None:
        bg_color = RGBColor(139, 21, 56)

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = bg_color

    # Titre centré en blanc
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Sous-titre
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(32)
        p2.font.color.rgb = RGBColor(255, 255, 255)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(20)

    return slide

def add_two_columns_slide(prs, title_text, left_title, left_content, right_title, right_content):
    """Ajoute une slide avec 2 colonnes de comparaison"""
    slide_layout = prs.slide_layouts[3]  # Layout deux colonnes
    slide = prs.slides.add_slide(slide_layout)

    # Titre principal
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(147, 22, 60)

    # Colonne gauche
    left_shape = slide.placeholders[1]
    left_tf = left_shape.text_frame
    left_tf.clear()

    p_left_title = left_tf.paragraphs[0]
    p_left_title.text = left_title
    p_left_title.font.size = Pt(24)
    p_left_title.font.bold = True
    p_left_title.font.color.rgb = RGBColor(59, 130, 246)  # Bleu
    p_left_title.space_after = Pt(15)

    for line in left_content:
        p = left_tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.space_before = Pt(6)

    # Colonne droite
    right_shape = slide.placeholders[2]
    right_tf = right_shape.text_frame
    right_tf.clear()

    p_right_title = right_tf.paragraphs[0]
    p_right_title.text = right_title
    p_right_title.font.size = Pt(24)
    p_right_title.font.bold = True
    p_right_title.font.color.rgb = RGBColor(16, 185, 129)  # Vert
    p_right_title.space_after = Pt(15)

    for line in right_content:
        p = right_tf.add_paragraph()
        p.text = line
        p.font.size = Pt(14)
        p.space_before = Pt(6)

    return slide

# Ouvrir la présentation existante
prs = Presentation('ALFlight_Presentation.pptx')

print(f"Presentation chargee. Nombre de slides actuelles: {len(prs.slides)}")

# ============================================
# SLIDE 9: Titre section "Avec ALFlight"
# ============================================
print("Ajout Slide 9: Titre section ALFlight...")
slide9 = add_transition_slide(
    prs,
    "Avec ALFlight",
    "La préparation de vol réinventée",
    bg_color=RGBColor(16, 185, 129)  # Vert success
)

# ============================================
# SLIDE 10: Vue d'ensemble 2 étapes
# ============================================
print("Ajout Slide 10: Vue d'ensemble 2 etapes...")
slide10_content = [
    "📱 ALFlight révolutionne la préparation de vol en 2 étapes :",
    "",
    "1️⃣ PREMIÈRE UTILISATION : Configuration initiale",
    "  ⏱️ Durée : 10-15 minutes (UNE SEULE FOIS)",
    "  • Créer votre profil pilote (licences, qualifications)",
    "  • Télécharger vos avions favoris depuis la bibliothèque",
    "  • Configurer vos préférences (unités, aérodromes habituels)",
    "",
    "2️⃣ CHAQUE VOL SUIVANT : Préparation instantanée",
    "  ⚡ Durée : 3 MINUTES CHRONO",
    "  • Saisir départ/arrivée → Tout le reste est AUTOMATIQUE",
    "  • Validation des calculs → Prêt à décoller !",
    "",
    "💡 Comparaison finale :",
    "  • Méthode manuelle : 3 heures À CHAQUE VOL",
    "  • ALFlight : 15 min la 1ère fois, puis 3 min par vol",
    "  ➜ Gain de temps : x60 plus rapide après configuration !",
]

add_title_content_slide(
    prs,
    "Comment Fonctionne ALFlight ?",
    slide10_content,
    subtitle="🎯 2 étapes simples pour une efficacité maximale"
)

# ============================================
# SLIDE 11: Première utilisation - Configuration
# ============================================
print("Ajout Slide 11: Premiere utilisation...")
slide11_content = [
    "👤 ÉTAPE 1 : Profil Pilote (3-5 min)",
    "  • Nom, prénom, date de naissance",
    "  • Licence (LAPL, PPL, CPL...) + numéro",
    "  • Qualifications (SEP, MEP, Night, IR...)",
    "  • Certificat médical (classe + date expiration)",
    "  ✅ Alertes automatiques pour renouvellements !",
    "",
    "✈️ ÉTAPE 2 : Télécharger Avions (2-3 min)",
    "  • Accéder bibliothèque intégrée (50+ modèles)",
    "  • Rechercher votre avion (DA40 NG, C172, PA28...)",
    "  • Télécharger : performances + abaques + limites",
    "  • Ajouter immatriculation (F-HSTR, F-GXYZ...)",
    "  ✅ Avion prêt avec TOUTES ses données techniques !",
    "",
    "⚙️ ÉTAPE 3 : Préférences Unités (2-3 min)",
    "  • Choisir preset rapide : EASA (Europe) ou FAA (USA)",
    "  • Ou personnaliser : distances (nm/km), vitesses (kt/km/h)...",
    "  • Configurer aérodromes favoris (base habituelle)",
    "  ✅ Application adaptée à VOS habitudes !",
    "",
    "📲 ÉTAPE 4 : Synchronisation (2-3 min)",
    "  • Créer compte Supabase (email + mot de passe)",
    "  • Sync automatique multi-appareils (phone + tablette + PC)",
    "  ✅ Vos données partout, tout le temps !",
]

add_title_content_slide(
    prs,
    "1️⃣ Première Utilisation : Configuration Initiale",
    slide11_content,
    subtitle="⏱️ Durée totale : 10-15 minutes (une seule fois dans votre vie de pilote !)"
)

# ============================================
# SLIDE 12: Transition vers usage quotidien
# ============================================
print("Ajout Slide 12: Transition usage quotidien...")
slide12 = add_transition_slide(
    prs,
    "Et Maintenant...",
    "Préparez CHAQUE vol en 3 minutes ⚡"
)

# ============================================
# SLIDE 13: Chaque vol - 3 minutes chrono
# ============================================
print("Ajout Slide 13: Chaque vol 3 minutes...")
slide13_content = [
    "🚀 MINUTE 1 : Saisie du Vol (30 secondes)",
    "  • Ouvrir ALFlight → Onglet 'Nouveau Vol'",
    "  • Saisir départ : LFST (Strasbourg)",
    "  • Saisir arrivée : LFPO (Paris-Orly)",
    "  • Sélectionner avion : F-HSTR (DA40 NG)",
    "  • Date/heure vol : Aujourd'hui 14h00",
    "  ✅ Clic sur 'Préparer Vol' → MAGIE !",
    "",
    "⚡ MINUTE 2 : ALFlight Fait TOUT Automatiquement (2 min)",
    "  ✅ Récupération METAR/TAF temps réel (LFST + LFPO)",
    "  ✅ Téléchargement NOTAM actifs (départ + arrivée + route)",
    "  ✅ Téléchargement VAC à jour (PDF automatique)",
    "  ✅ Calcul route optimale + waypoints GPS",
    "  ✅ Calcul vents en altitude (FL50, FL75, FL100...)",
    "  ✅ Calcul temps de vol et carburant nécessaire",
    "  ✅ Calcul performances décollage (PA, DA, distance, marges)",
    "  ✅ Calcul performances atterrissage (avec dégagement)",
    "  ✅ Calcul Weight & Balance automatique",
    "  ✅ Génération plan de vol OLIVIA pré-rempli",
    "  ✅ Analyse espaces aériens traversés (TMA, CTR...)",
    "  ✅ Export PDF complet prêt à imprimer",
    "",
    "👀 MINUTE 3 : Validation Pilote (30 secondes)",
    "  • Vérifier résultats affichés (distances, temps, fuel)",
    "  • Lire NOTAM critiques surlignés en rouge",
    "  • Valider météo VMC/IMC acceptable",
    "  • Confirmer performances SAFE (verdicts verts)",
    "  ✅ C'EST PRÊT ! Direction l'avion 🛫",
]

add_title_content_slide(
    prs,
    "2️⃣ Chaque Vol : Préparation en 3 Minutes ⚡",
    slide13_content,
    subtitle="⏱️ Chronomètre : De l'ouverture de l'app au décollage en 180 secondes"
)

# ============================================
# SLIDE 14: Comparaison avant/après
# ============================================
print("Ajout Slide 14: Comparaison avant/apres...")
left_content = [
    "⏱️ TEMPS : 3 heures par vol",
    "📚 15-20 sites web à consulter",
    "🧮 50+ calculs manuels",
    "📄 Documents papier éparpillés",
    "❌ Risques d'erreurs élevés",
    "😤 Frustration et stress",
    "🔄 Tout recommencer si changement",
    "💸 Ratio 2:1 (prépa:vol)",
    "",
    "👎 Décourage les vols spontanés",
]

right_content = [
    "⚡ TEMPS : 3 minutes par vol",
    "📱 1 seule application tout-en-un",
    "🤖 100% automatisé (0 calcul manuel)",
    "📲 Tout centralisé dans l'app",
    "✅ Calculs vérifiés et précis",
    "😌 Sérénité et confiance",
    "🔄 Mise à jour instantanée (1 clic)",
    "💸 Ratio 1:30 (prépa:vol)",
    "",
    "👍 Encourage la pratique régulière",
]

slide14 = add_two_columns_slide(
    prs,
    "Avant / Après : Le Changement Radical",
    "❌ AVANT (Sans ALFlight)",
    left_content,
    "✅ APRÈS (Avec ALFlight)",
    right_content
)

# ============================================
# SLIDE 15: Le vrai gain
# ============================================
print("Ajout Slide 15: Le vrai gain...")
slide15_content = [
    "💡 Le VRAI gain d'ALFlight n'est pas que le temps...",
    "",
    "🧠 CHARGE MENTALE RÉDUITE DE 95%",
    "  • Fini le stress de l'oubli d'un NOTAM critique",
    "  • Fini la peur d'une erreur de calcul dans les performances",
    "  • Fini les doutes sur la validité des documents",
    "  ➜ CONFIANCE TOTALE dans votre préparation",
    "",
    "🛫 PLUS DE VOLS, PLUS SOUVENT",
    "  • Vol spontané le weekend ? OUI en 3 minutes !",
    "  • Météo change ? Pas de problème, recalcul instantané",
    "  • Passager supplémentaire ? 1 clic, W&B refait",
    "  ➜ LIBERTÉ de voler quand vous voulez",
    "",
    "📈 PROGRESSION ACCÉLÉRÉE",
    "  • Plus de temps en vol, moins en paperasse",
    "  • Statistiques automatiques (heures PIC, SEP, nuit...)",
    "  • Alertes revalidations licences (jamais périmé)",
    "  ➜ CARRIÈRE de pilote optimisée",
    "",
    "💰 RETOUR SUR INVESTISSEMENT",
    "  • Abonnement : 9,99€/mois",
    "  • Temps gagné par vol : 3h → 0,05h = 2h55 économisées",
    "  • Sur 1 an (20 vols) : 58 heures récupérées",
    "  ➜ Soit 58h × (coût location avion) en vols supplémentaires !",
]

add_title_content_slide(
    prs,
    "🎯 Le Véritable Impact d'ALFlight",
    slide15_content,
    subtitle="Au-delà du temps : Confiance, Liberté, Progression"
)

# ============================================
# SLIDE 16: Call to action
# ============================================
print("Ajout Slide 16: Call to action...")
slide16 = add_transition_slide(
    prs,
    "Prêt à Révolutionner",
    "Votre Pratique du Vol ? 🚀",
    bg_color=RGBColor(147, 22, 60)  # Retour bordeaux ALFlight
)

print(f"\nSlides ajoutees avec succes ! Total slides : {len(prs.slides)}")

# Sauvegarder la présentation avec nouveau nom
output_path = 'ALFlight_Presentation_Complete.pptx'
prs.save(output_path)
print(f"Presentation mise a jour : {output_path}")
print("\nContenu ajoute :")
print("  - Slide 9 : Titre section 'Avec ALFlight'")
print("  - Slide 10 : Vue d'ensemble 2 etapes (15 min + 3 min)")
print("  - Slide 11 : Premiere utilisation - Configuration (10-15 min)")
print("  - Slide 12 : Transition vers usage quotidien")
print("  - Slide 13 : Chaque vol en 3 minutes chrono")
print("  - Slide 14 : Comparaison avant/apres (2 colonnes)")
print("  - Slide 15 : Le vrai gain (confiance, liberte, progression)")
print("  - Slide 16 : Call to action")
print("\nTotal : 8 nouvelles slides workflow ALFlight")
print("Presentation complete : 20 + 8 = 28 slides")
