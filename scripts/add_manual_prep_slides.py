"""
Script pour ajouter des slides sur la préparation de vol manuelle (sans ALFlight)
Crée un contraste brutal entre la méthode traditionnelle et l'app
"""

import subprocess
import sys

# Installer python-pptx si nécessaire
try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
except ImportError:
    print("Installation de python-pptx...")
    subprocess.check_call([sys.executable, "-m", "pip", "install", "python-pptx"])
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor

def add_title_content_slide(prs, title_text, content_lines, subtitle=None):
    """Ajoute une slide avec titre et contenu bullet points"""
    slide_layout = prs.slide_layouts[1]  # Layout avec titre et contenu
    slide = prs.slides.add_slide(slide_layout)

    # Titre
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(40)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(147, 22, 60)  # Bordeaux ALFlight

    # Contenu
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()

    # Sous-titre optionnel
    if subtitle:
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.italic = True
        p.font.color.rgb = RGBColor(107, 15, 43)
        p.space_after = Pt(20)

    # Lignes de contenu
    for i, line in enumerate(content_lines):
        if i == 0 and not subtitle:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.text = line
        p.level = 0 if not line.startswith("  ") else 1
        p.font.size = Pt(16) if p.level == 0 else Pt(14)
        p.space_before = Pt(8) if p.level == 0 else Pt(4)

        # Couleur rouge pour les warnings
        if "⚠️" in line or "❌" in line or "ATTENTION" in line:
            p.font.color.rgb = RGBColor(239, 68, 68)
            p.font.bold = True

    return slide

def add_transition_slide(prs, title_text, subtitle_text):
    """Ajoute une slide de transition dramatique"""
    slide_layout = prs.slide_layouts[5]  # Layout vide
    slide = prs.slides.add_slide(slide_layout)

    # Fond bordeaux
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(139, 21, 56)

    # Titre centré en blanc
    left = Inches(1)
    top = Inches(2.5)
    width = Inches(8)
    height = Inches(2)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)
    p.alignment = PP_ALIGN.CENTER

    # Sous-titre
    if subtitle_text:
        p2 = tf.add_paragraph()
        p2.text = subtitle_text
        p2.font.size = Pt(32)
        p2.font.color.rgb = RGBColor(255, 255, 255)
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(20)

    return slide

# Ouvrir la présentation existante
prs = Presentation('ALFlight_Presentation.pptx')

print(f"Présentation chargée. Nombre de slides actuelles: {len(prs.slides)}")

# ============================================
# SLIDE 1: Titre choc
# ============================================
print("Ajout Slide 1: Titre choc...")
slide1 = add_transition_slide(
    prs,
    "⏱️ Préparer un vol SANS ALFlight...",
    "Bienvenue dans le cauchemar du pilote moderne"
)

# ============================================
# SLIDE 2: Recherche et collecte des documents
# ============================================
print("Ajout Slide 2: Recherche documents...")
slide2_content = [
    "📄 Rechercher manuellement chaque NOTAM (15-20 min)",
    "  • Ouvrir le site SIA (Aeroweb) → Naviguer dans les menus",
    "  • Chercher les NOTAM aérodrome de départ (LFST)",
    "  • Chercher les NOTAM aérodrome d'arrivée (LFPO)",
    "  • Chercher les NOTAM en route (zones traversées)",
    "  • Lire et trier manuellement les NOTAM pertinents",
    "  • Noter sur papier ceux qui impactent le vol",
    "",
    "🌦️ Récupérer les METAR/TAF (5-10 min)",
    "  • Ouvrir un autre site météo (Météo France, AWC...)",
    "  • Chercher METAR départ, arrivée, dégagements",
    "  • Chercher TAF pour l'heure prévue",
    "  • Décoder manuellement les METAR/TAF (groupes nuages, vents...)",
    "  • Évaluer si les conditions sont VMC/IMC",
    "",
    "🗺️ Télécharger les cartes VAC (10-15 min)",
    "  • Télécharger PDF VAC départ (SIA)",
    "  • Télécharger PDF VAC arrivée",
    "  • Vérifier dates de validité (périmées ?)",
    "  • Imprimer les VAC (ou les avoir sur tablette séparée)",
    "",
    "📚 Sortir le manuel de vol avion (5 min)",
    "  • Chercher dans la sacoche/cockpit le POH papier",
    "  • S'assurer d'avoir la bonne version (DA40 NG, C172, PA28...)",
]

add_title_content_slide(
    prs,
    "Étape 1/5 : Recherche et Collecte des Documents",
    slide2_content,
    subtitle="⏱️ Durée moyenne : 45-60 minutes (si tout va bien...)"
)

# ============================================
# SLIDE 3: Planification de la route
# ============================================
print("Ajout Slide 3: Planification route...")
slide3_content = [
    "🗺️ Tracer la route manuellement (20-30 min)",
    "  • Sortir la carte OACI papier 1/500 000",
    "  • Tracer au crayon la route entre LFST et LFPO",
    "  • Choisir des waypoints visuels (villes, lacs, routes...)",
    "  • Mesurer les distances avec une règle graduée",
    "  • Mesurer les caps magnétiques avec un rapporteur",
    "  • Corriger les caps pour la déclinaison magnétique (+1° en France)",
    "",
    "⛰️ Vérifier les altitudes minimales (10-15 min)",
    "  • Identifier le relief sur la carte (courbes de niveau)",
    "  • Trouver le point culminant de chaque segment",
    "  • Ajouter marge de sécurité 1000 ft (VFR) ou 2000 ft (montagne)",
    "  • Vérifier les planchers des espaces aériens (TMA, CTR...)",
    "",
    "✈️ Calculer temps de vol et carburant (15-20 min)",
    "  • Estimer la vitesse sol (TAS + vent)",
    "  • ❌ Problème : Vent sur route ? Faut chercher les vents en altitude...",
    "  • Ouvrir ENCORE un autre site (Windy, Meteoblue...)",
    "  • Calculer temps de vol segment par segment : T = D / GS",
    "  • Calculer consommation : Fuel = Time × Flow (chercher flow dans POH...)",
    "  • Ajouter réserves réglementaires (30 min jour, 45 min nuit)",
]

add_title_content_slide(
    prs,
    "Étape 2/5 : Planification de la Route",
    slide3_content,
    subtitle="⏱️ Durée moyenne : 45-60 minutes"
)

# ============================================
# SLIDE 4: Calculs de performances
# ============================================
print("Ajout Slide 4: Calculs performances...")
slide4_content = [
    "📊 Performances décollage (15-20 min)",
    "  • Chercher la température actuelle dans le METAR (décodage manuel)",
    "  • Calculer Altitude Pression : PA = Elev + (1013-QNH) × 27",
    "  • ❌ Problème : Besoin d'une calculatrice...",
    "  • Calculer Altitude Densité : DA = PA + 120 × (OAT - ISA)",
    "  • Ouvrir le POH à la section 'Takeoff Performance'",
    "  • Trouver l'abaque correspondant (piste herbe ou dur ?)",
    "  • Lire graphiquement la distance de décollage (imprécis ±10%)",
    "  • Appliquer facteurs correctifs manuels :",
    "    - Vent de face/arrière : ±10% par 10kt",
    "    - Pente piste : +10% par 1% de pente montante",
    "    - Herbe : +20% de la distance",
    "    - Piste humide : +15%",
    "  • Ajouter marge sécurité 50% (réglementaire)",
    "  • ⚠️ ATTENTION : Oublier UN SEUL facteur = sous-estimation dangereuse !",
    "",
    "📊 Performances atterrissage (15 min)",
    "  • Chercher METAR destination (vent, température...)",
    "  • Refaire TOUS les calculs PA/DA pour l'arrivée",
    "  • Lire abaque atterrissage dans POH",
    "  • Appliquer facteurs correctifs (vent, pente, surface...)",
    "  • Vérifier longueur piste suffisante avec marges",
]

add_title_content_slide(
    prs,
    "Étape 3/5 : Calculs de Performances",
    slide4_content,
    subtitle="⏱️ Durée moyenne : 30-45 minutes (avec erreurs fréquentes...)"
)

# ============================================
# SLIDE 5: Weight & Balance et check final
# ============================================
print("Ajout Slide 5: W&B et checks...")
slide5_content = [
    "⚖️ Masse et Centrage (15-20 min)",
    "  • Chercher masse à vide avion (Carnet de Route)",
    "  • Calculer masse équipage : 2 personnes × 77 kg (forfait EASA)",
    "  • Calculer masse bagages (estimation...)",
    "  • Calculer masse carburant : Fuel(L) × 0.72 kg/L",
    "  • Additionner pour avoir masse totale",
    "  • Vérifier < MTOW (Maximum Takeoff Weight)",
    "  • Calculer centrage (moments) :",
    "    - Moment = Masse × Bras de levier",
    "    - Chercher les bras de levier dans le POH (tableau complexe)",
    "    - Calculer moment total et CG position",
    "  • Reporter sur l'enveloppe de centrage POH",
    "  • ❌ Refaire si hors enveloppe (déplacer bagages, enlever carburant...)",
    "",
    "📋 Remplir le plan de vol (10-15 min)",
    "  • Remplir formulaire OLIVIA (plan de vol DGAC)",
    "  • Saisir manuellement : aérodromes, route, heures, autonomie...",
    "  • Déposer en ligne (ou par téléphone si site en panne...)",
    "",
    "✅ Check final liste documents bord (5 min)",
    "  • Vérifier tous les documents imprimés/téléchargés",
    "  • S'assurer de ne rien avoir oublié (facile d'oublier un NOTAM...)",
]

add_title_content_slide(
    prs,
    "Étape 4/5 : Weight & Balance et Check Final",
    slide5_content,
    subtitle="⏱️ Durée moyenne : 30-40 minutes"
)

# ============================================
# SLIDE 6: Le cauchemar des changements
# ============================================
print("Ajout Slide 6: Cauchemar changements...")
slide6_content = [
    "❌ Changement météo ? TOUT RECOMMENCER !",
    "  • Nouveau METAR → Nouvelles températures",
    "  • Recalculer PA, DA, performances décollage/atterrissage",
    "  • Nouveaux vents → Recalculer vitesses sol, temps, carburant",
    "  • Conditions détériorées → Chercher un aérodrome de dégagement",
    "  • Refaire les calculs pour le dégagement (METAR, VAC, perfs...)",
    "",
    "❌ Passager supplémentaire ? TOUT RECOMMENCER !",
    "  • Recalculer masse et centrage complet",
    "  • Nouvelles performances décollage (masse augmentée)",
    "  • Recalculer consommation carburant (avion plus lourd)",
    "  • Potentiellement : retirer des bagages ou du carburant",
    "  • Si hors enveloppe centrage → Reorganiser chargement",
    "",
    "❌ Changement de piste ? Recalculer !",
    "  • Nouvelle orientation piste → Recalculer composantes vent",
    "  • Nouvelle longueur piste → Revérifier performances",
    "",
    "❌ Retard au départ ? Mise à jour !",
    "  • Nouveaux horaires → Nouveau TAF",
    "  • Nouveau METAR → Températures/vents différents",
    "  • Plan de vol à modifier (OLIVIA)",
]

add_title_content_slide(
    prs,
    "Étape 5/5 : Le Cauchemar des Changements",
    slide6_content,
    subtitle="⚠️ CHAQUE modification = Recommencer partiellement ou totalement"
)

# ============================================
# SLIDE 7: Le bilan catastrophique
# ============================================
print("Ajout Slide 7: Bilan...")
slide7_content = [
    "⏱️ TEMPS TOTAL : 2h30 à 3h30 de préparation",
    "  • Pour UN SEUL vol de navigation",
    "  • Sans compter les erreurs et oublis...",
    "",
    "🧠 CHARGE MENTALE : Maximale",
    "  • 15-20 sites web différents à consulter",
    "  • 50+ calculs manuels (risques d'erreur élevés)",
    "  • Documents papier éparpillés partout",
    "  • Impossible de tout garder en tête",
    "",
    "⚠️ RISQUES D'ERREURS : Très élevés",
    "  • Oubli d'un NOTAM critique (piste fermée)",
    "  • Erreur de calcul dans les performances (dangereux !)",
    "  • Mauvaise lecture d'un abaque (±15% d'erreur)",
    "  • Confusion entre unités (ft/m, kt/km/h, L/gal...)",
    "  • Document périmé non détecté",
    "",
    "😤 FRUSTRATION : Maximum",
    "  • Répétitif, chronophage, source de stress",
    "  • Décourage les vols spontanés",
    "  • Vol de 1h30 = 3h de préparation (ratio 2:1 !)",
    "",
    "💸 COÛT OPPORTUNITÉ : Énorme",
    "  • 3 heures = Temps qu'on pourrait passer À VOLER",
]

add_title_content_slide(
    prs,
    "📊 Le Bilan Catastrophique",
    slide7_content,
    subtitle="🚨 La réalité de CHAQUE préparation de vol sans aide numérique"
)

# ============================================
# SLIDE 8: Transition dramatique
# ============================================
print("Ajout Slide 8: Transition...")
slide8 = add_transition_slide(
    prs,
    "ET SI...",
    "Tout cela pouvait prendre 15 minutes ?"
)

print(f"\nSlides ajoutées avec succès ! Total slides : {len(prs.slides)}")

# Sauvegarder la présentation
output_path = 'ALFlight_Presentation_Updated.pptx'
prs.save(output_path)
print(f"✅ Présentation sauvegardée : {output_path}")
print("\nContenu ajouté :")
print("  - Slide 1 : Titre choc")
print("  - Slide 2 : Recherche documents (45-60 min)")
print("  - Slide 3 : Planification route (45-60 min)")
print("  - Slide 4 : Calculs performances (30-45 min)")
print("  - Slide 5 : W&B et check final (30-40 min)")
print("  - Slide 6 : Cauchemar des changements")
print("  - Slide 7 : Bilan catastrophique")
print("  - Slide 8 : Transition 'ET SI...'")
print("\n🎯 Contraste brutal créé : 3h de galère VS 15 min avec ALFlight !")
